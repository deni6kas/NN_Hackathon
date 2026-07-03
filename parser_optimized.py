# -*- coding: utf-8 -*-
"""
Оптимизированная версия parser.py.

Отличия от оригинала (функционал извлечения знаний из текста сохранён):
  1. Параллельные LLM-вызовы (ThreadPoolExecutor, ANALYSIS_WORKERS потоков)
     вместо строго последовательных.
  2. Крупные чанки текста (MAX_CHUNK_CHARS) вместо "1 страница = 1 вызов API" —
     в разы меньше запросов к Yandex GPT.
  3. Убрана обработка изображений (по заданию требуется NLP-пайплайн по текстам;
     каждая картинка стоила отдельного LLM-вызова).
  4. Накопление строк в списках вместо pd.concat на каждую строку
     (у оригинала квадратичная сложность на больших объёмах).
  5. Ретраи с экспоненциальной паузой при ошибках API.
  6. Промежуточное сохранение таблиц после каждого батча из CHECKPOINT_EVERY файлов.

Все таблицы пишутся в папку скрипта (parser_optimized/), чтобы не конфликтовать
с запущенным оригинальным парсером. Данные читаются из ../raw_data.
"""

import os
import json
import time
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
from datetime import datetime
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from openai import OpenAI

# Секреты берутся из переменных окружения (не хардкодить в публичном репозитории):
#   setx YANDEX_API_KEY   "ваш-ключ"
#   setx YANDEX_FOLDER_ID "ваш-folder-id"
API_KEY = os.environ["YANDEX_API_KEY"]
FOLDER_ID = os.environ["YANDEX_FOLDER_ID"]

BASE_DIR = Path(__file__).resolve().parent          # папка parser_optimized/
DATA_DIR = BASE_DIR.parent / "raw_data"             # исходные документы
CROSS_DOC_DB_FILE = BASE_DIR / "cross_document_links.json"

# ===== ПАРАМЕТРЫ ОПТИМИЗАЦИИ =====
ANALYSIS_WORKERS = 8      # параллельных запросов (квота Yandex: 10 одновременных сессий)
MAX_CHUNK_CHARS = 4000    # размер чанка (меньше — чтобы JSON-ответ не обрезался по лимиту вывода)
CHECKPOINT_EVERY = 10     # сохранять таблицы после каждых N обработанных файлов
MAX_RETRIES = 5           # повторы при ошибке API

client = OpenAI(
    api_key=API_KEY,
    base_url="https://llm.api.cloud.yandex.net/v1"
)
MODEL_URI = f"gpt://{FOLDER_ID}/yandexgpt/latest"

SYSTEM_PROMPT = """
Ты — интеллектуальный ассистент R&D лаборатории в горно-металлургической отрасли.
Твоя задача — проанализировать текст и извлечь из него ВСЮ ИНФОРМАЦИЮ строго в формате JSON.

ТИПЫ СУЩНОСТЕЙ: Material, Process, Equipment, Property, Experiment, Publication, Expert, Facility.
ТИПЫ СВЯЗЕЙ: uses_material, operates_at_condition, produces_output, described_in, validated_by, contradicts, researched_by, optimal_parameter.
ДОМЕНЫ: гидрометаллургия, пирометаллургия, экология, переработка отходов.
ГЕОГРАФИЧЕСКИЙ ОХВАТ: Russia|Worldwide|[страны].
УРОВЕНЬ ВЕРИФИКАЦИИ: Confirmed|Contradicted|Unverified|Partially_Confirmed.

Выведи ответ СТРОГО в этом JSON формате (ничего, кроме JSON):
{
  "nodes": [
    {
      "entity_type": "Тип сущности",
      "name": "Название",
      "language": "ru|en",
      "synonyms": "Синонимы если есть",
      "description": "Описание",
      "domain": "Домен (гидрометаллургия, пирометаллургия, экология, переработка отходов)",
      "is_canonical": true
    }
  ],
  "edges": [
    {
      "source_name": "ИмяУзла1",
      "target_name": "ИмяУзла2",
      "relationship_type": "ТипСвязи",
      "conditions": "Условия применения если есть (pH: 7-8, концентрация: 200-300 мг/л)",
      "geographic_scope": "Russia|Worldwide|Страны",
      "applicable_range": "Диапазоны численных значений если есть",
      "verification_level": "Уровень подтверждения"
    }
  ],
  "parameters": [
    {
      "entity_name": "Имя сущности, к которой относится параметр",
      "parameter_name": "Имя параметра (концентрация, температура, скорость потока, pH)",
      "unit": "Единица измерения (мг/л, °C, м³/ч)",
      "value": "Точное значение если известно",
      "min_value": "Минимальное значение",
      "max_value": "Максимальное значение",
      "context": "При каких условиях этот параметр актуален"
    }
  ],
  "experts": [
    {
      "name": "ФИО или организация",
      "organization": "Организация",
      "country": "Страна",
      "expertise_topics": "Области экспертизы (процессы/материалы которые изучают)",
      "language": "ru|en"
    }
  ],
  "contradictions": [
    {
      "description": "Описание противоречия",
      "topic": "О чем противоречие",
      "different_claims": "Разные утверждения"
    }
  ]
}
"""

# ===== КОЛОНКИ ТАБЛИЦ (как в оригинале) =====
NODES_COLS = ["node_id", "entity_type", "name", "language", "synonyms", "description",
              "domain", "source_document", "created_date", "updated_date", "is_canonical"]
EDGES_COLS = ["source_id", "target_id", "relationship_type", "weight", "bidirectional",
              "conditions", "geographic_scope", "applicable_range", "verification_level",
              "source_document", "confidence_level", "extraction_date"]
PARAMS_COLS = ["parameter_id", "node_id", "parameter_name", "unit", "value", "min_value",
               "max_value", "context", "confidence_level", "source_document", "extraction_date"]
EXPERTS_COLS = ["expert_id", "name", "organization", "country", "expertise_node_ids",
                "contact_email", "affiliation_year", "language", "verification_status"]
CONTRADICTIONS_COLS = ["contradiction_id", "node_id_1", "node_id_2", "relationship_type",
                       "description", "conflicting_sources", "resolution_status"]
DOC_META_COLS = ["doc_id", "doc_name", "doc_type", "authors", "publication_year",
                 "publication_venue", "country_origin", "language", "relevance_tags",
                 "access_level", "import_date", "verification_status"]


def get_ai_analysis(text_chunk):
    """Анализ текста через Yandex GPT (с ретраями)"""
    for attempt in range(MAX_RETRIES):
        try:
            response = client.chat.completions.create(
                model=MODEL_URI,
                messages=[
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": f"Проанализируй следующий текст:\n\n{text_chunk}"}
                ],
                temperature=0.1
            )
            raw_text = response.choices[0].message.content
            clean_json = raw_text.strip().replace("```json", "").replace("```", "").strip()
            return json.loads(clean_json)
        except Exception as e:
            msg = str(e)
            if attempt >= MAX_RETRIES - 1:
                print(f"  ✗ Ошибка ИИ (текст, {MAX_RETRIES} попыток): {e}")
            elif "429" in msg or "rate_limit" in msg:
                # квота одновременных запросов — ждём подольше, запрос не потерян
                time.sleep(5 * (attempt + 1))
            else:
                time.sleep(2 ** attempt)
    return None


def merge_pieces(pieces, max_chars=MAX_CHUNK_CHARS):
    """Склейка мелких кусков текста (страниц/абзацев/слайдов) в крупные чанки"""
    chunks = []
    current = ""
    for piece in pieces:
        piece = piece.strip()
        if not piece:
            continue
        if current and len(current) + len(piece) + 1 > max_chars:
            chunks.append(current)
            current = piece
        else:
            current = f"{current}\n{piece}" if current else piece
    if current:
        chunks.append(current)
    return chunks


def extract_text_chunks(file_path, ext):
    """Извлечение текста из документа и разбивка на крупные чанки"""
    pieces = []
    if ext == '.pdf':
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page_text = doc[page_num].get_text().strip()
            if page_text:
                pieces.append(page_text)
        doc.close()
    elif ext in ('.docx', '.docm'):
        doc = Document(file_path)
        pieces = [p.text for p in doc.paragraphs if p.text.strip()]
    elif ext == '.pptx':
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = [shape.text.strip() for shape in slide.shapes
                          if hasattr(shape, "text") and shape.text.strip()]
            if slide_text:
                pieces.append("\n".join(slide_text))
    return merge_pieces(pieces)


class KnowledgeStore:
    """Накопитель результатов: списки вместо pd.concat на каждую строку"""

    def __init__(self):
        self.nodes = self._load("nodes.csv", NODES_COLS)
        self.edges = self._load("edges.csv", EDGES_COLS)
        self.parameters = self._load("parameters.csv", PARAMS_COLS)
        self.experts = self._load("experts.csv", EXPERTS_COLS)
        self.contradictions = self._load("contradictions.csv", CONTRADICTIONS_COLS)
        self.document_metadata = self._load("document_metadata.csv", DOC_META_COLS)

        self.node_id_by_name = {str(n["name"]).lower(): n["node_id"] for n in self.nodes}
        self.node_counter = len(self.nodes) + 1
        self.parameter_counter = len(self.parameters) + 1
        self.expert_counter = len(self.experts) + 1
        self.contradiction_counter = len(self.contradictions) + 1
        self.doc_counter = len(self.document_metadata) + 1

    @staticmethod
    def _load(filename, columns):
        path = BASE_DIR / filename
        if path.exists():
            return pd.read_csv(path).to_dict("records")
        return []

    def add_document_metadata(self, file_name, doc_type):
        self.document_metadata.append({
            "doc_id": f"DOC_{self.doc_counter}",
            "doc_name": file_name,
            "doc_type": doc_type,
            "authors": "",
            "publication_year": "",
            "publication_venue": "",
            "country_origin": "Russia",
            "language": "ru",
            "relevance_tags": "",
            "access_level": "internal",
            "import_date": datetime.now().strftime("%Y-%m-%d"),
            "verification_status": "Unverified"
        })
        self.doc_counter += 1

    def merge_analysis(self, result, file_name):
        """Слияние результата LLM-анализа одного чанка (вызывается из главного потока)"""
        if not result:
            return

        today = datetime.now().strftime("%Y-%m-%d")
        local_name_to_id = {}

        # Узлы
        for node in result.get("nodes", []):
            if not isinstance(node, dict) or not node.get("name"):
                continue
            name_lower = str(node["name"]).lower()
            if name_lower not in self.node_id_by_name:
                prefix = str(node.get("entity_type", "OTH"))[:3].upper()
                node_id = f"{prefix}_{self.node_counter}"
                self.nodes.append({
                    "node_id": node_id,
                    "entity_type": node.get("entity_type", "Other"),
                    "name": node["name"],
                    "language": node.get("language", "ru"),
                    "synonyms": node.get("synonyms", ""),
                    "description": node.get("description", f"Из {file_name}"),
                    "domain": node.get("domain", ""),
                    "source_document": file_name,
                    "created_date": today,
                    "updated_date": today,
                    "is_canonical": node.get("is_canonical", True)
                })
                self.node_id_by_name[name_lower] = node_id
                self.node_counter += 1
            local_name_to_id[name_lower] = self.node_id_by_name[name_lower]

        # Связи
        for edge in result.get("edges", []):
            if not isinstance(edge, dict):
                continue
            src_lower = str(edge.get("source_name", "")).lower()
            tgt_lower = str(edge.get("target_name", "")).lower()
            if src_lower in local_name_to_id and tgt_lower in local_name_to_id:
                self.edges.append({
                    "source_id": local_name_to_id[src_lower],
                    "target_id": local_name_to_id[tgt_lower],
                    "relationship_type": edge.get("relationship_type", ""),
                    "weight": 0.5,
                    "bidirectional": False,
                    "conditions": edge.get("conditions", ""),
                    "geographic_scope": edge.get("geographic_scope", "Worldwide"),
                    "applicable_range": edge.get("applicable_range", ""),
                    "verification_level": edge.get("verification_level", "Unverified"),
                    "source_document": file_name,
                    "confidence_level": "High",
                    "extraction_date": today
                })

        # Параметры
        for param in result.get("parameters", []):
            if not isinstance(param, dict):
                continue
            entity_name = str(param.get("entity_name", "")).lower()
            if entity_name in local_name_to_id:
                self.parameters.append({
                    "parameter_id": f"PAR_{self.parameter_counter}",
                    "node_id": local_name_to_id[entity_name],
                    "parameter_name": param.get("parameter_name", ""),
                    "unit": param.get("unit", ""),
                    "value": param.get("value", ""),
                    "min_value": param.get("min_value", ""),
                    "max_value": param.get("max_value", ""),
                    "context": param.get("context", ""),
                    "confidence_level": "High",
                    "source_document": file_name,
                    "extraction_date": today
                })
                self.parameter_counter += 1

        # Эксперты
        for expert in result.get("experts", []):
            if not isinstance(expert, dict):
                continue
            self.experts.append({
                "expert_id": f"EXP_{self.expert_counter}",
                "name": expert.get("name", ""),
                "organization": expert.get("organization", ""),
                "country": expert.get("country", ""),
                "expertise_node_ids": expert.get("expertise_topics", ""),
                "contact_email": "",
                "affiliation_year": "",
                "language": expert.get("language", "ru"),
                "verification_status": "Unverified"
            })
            self.expert_counter += 1

        # Противоречия
        for contradiction in result.get("contradictions", []):
            if not isinstance(contradiction, dict):
                continue
            self.contradictions.append({
                "contradiction_id": f"CONTR_{self.contradiction_counter}",
                "node_id_1": "",
                "node_id_2": "",
                "relationship_type": "contradicts",
                "description": contradiction.get("description", ""),
                "conflicting_sources": file_name,
                "resolution_status": "open"
            })
            self.contradiction_counter += 1

    def save(self):
        """Сохранение всех таблиц в CSV (в папку parser_optimized/)"""
        tables = [
            ("nodes.csv", self.nodes, NODES_COLS),
            ("edges.csv", self.edges, EDGES_COLS),
            ("parameters.csv", self.parameters, PARAMS_COLS),
            ("experts.csv", self.experts, EXPERTS_COLS),
            ("contradictions.csv", self.contradictions, CONTRADICTIONS_COLS),
            ("document_metadata.csv", self.document_metadata, DOC_META_COLS),
            ("facilities.csv", [], []),
            ("keywords_synonyms.csv", [], []),
        ]
        for filename, rows, columns in tables:
            df = pd.DataFrame(rows, columns=columns if columns else None)
            df.to_csv(BASE_DIR / filename, index=False, encoding="utf-8-sig")


def deduplicate_edges(edges_df):
    """Дедупликация и объединение идентичных связей (как в оригинале)"""
    if len(edges_df) == 0:
        return edges_df

    if 'bidirectional' not in edges_df.columns:
        edges_df['bidirectional'] = False

    if 'weight' not in edges_df.columns:
        confidence_map = {'High': 0.95, 'Medium': 0.7, 'Low': 0.4}
        edges_df['weight'] = edges_df['confidence_level'].map(
            lambda x: confidence_map.get(x, 0.5)
        )

    grouped = edges_df.groupby(['source_id', 'target_id', 'relationship_type'], as_index=False).agg({
        'confidence_level': lambda x: x.iloc[0],
        'weight': 'max',
        'source_document': lambda x: '|'.join(x.astype(str).unique()),
        'extraction_date': 'max',
        'value_range': lambda x: '|'.join(x.dropna().unique()) if x.notna().any() else ''
    } if 'value_range' in edges_df.columns else {
        'confidence_level': lambda x: x.iloc[0],
        'weight': 'max',
        'source_document': lambda x: '|'.join(x.astype(str).unique()),
        'extraction_date': 'max'
    })

    # Двусторонние связи: векторно, по множеству перевёрнутых пар
    pairs = set(zip(grouped['source_id'], grouped['target_id'], grouped['relationship_type']))
    grouped['bidirectional'] = [
        (t, s, r) in pairs for s, t, r in
        zip(grouped['source_id'], grouped['target_id'], grouped['relationship_type'])
    ]

    return grouped


def build_cross_document_links(nodes):
    """Построение связей между документами на основе общих сущностей (как в оригинале)"""
    doc_entities = {}
    for node in nodes:
        doc = node.get("source_document")
        if doc and pd.notna(doc):
            doc_entities.setdefault(doc, set()).add(str(node["name"]).lower())

    cross_doc_links = []
    doc_names = list(doc_entities.keys())
    for i, doc1 in enumerate(doc_names):
        for doc2 in doc_names[i + 1:]:
            common_entities = doc_entities[doc1] & doc_entities[doc2]
            if common_entities:
                cross_doc_links.append({
                    "document_1": doc1,
                    "document_2": doc2,
                    "common_entities": list(common_entities),
                    "count": len(common_entities),
                    "link_type": "common_entities",
                    "confidence": "high" if len(common_entities) > 2 else "medium"
                })

    if cross_doc_links:
        with open(CROSS_DOC_DB_FILE, 'w', encoding='utf-8') as f:
            json.dump(cross_doc_links, f, ensure_ascii=False, indent=2)
        print(f"\n🔗 Найдено {len(cross_doc_links)} кросс-документных связей")

    return cross_doc_links


def main():
    start_time = time.time()
    store = KnowledgeStore()

    supported_extensions = ('.pdf', '.docx', '.docm', '.pptx')
    file_tasks = []
    for root, dirs, files in os.walk(DATA_DIR):
        for file in files:
            if file.lower().endswith(supported_extensions):
                file_tasks.append((os.path.join(root, file), file))

    if not file_tasks:
        print(f"Файлы не найдены в директории {DATA_DIR} или её подпапках.")
        return

    print(f"Всего в структуре найдено {len(file_tasks)} файлов.")

    # Возобновление: пропускаем файлы, уже сохранённые в document_metadata на чекпоинтах
    already_processed = {m["doc_name"] for m in store.document_metadata}
    if already_processed:
        file_tasks = [(p, n) for p, n in file_tasks if n not in already_processed]
        print(f"Уже обработано ранее (пропускаю): {len(already_processed)}, осталось: {len(file_tasks)}")

    total = len(file_tasks)
    if not file_tasks:
        print("Все файлы уже обработаны.")
        return
    print(f"Параллельных LLM-запросов: {ANALYSIS_WORKERS}, размер чанка: {MAX_CHUNK_CHARS} символов\n")

    all_processed_documents = []
    processed = 0

    # Обработка батчами по CHECKPOINT_EVERY файлов:
    # извлечение текста -> параллельный LLM-анализ чанков -> слияние -> checkpoint
    with ThreadPoolExecutor(max_workers=ANALYSIS_WORKERS) as executor:
        for batch_start in range(0, total, CHECKPOINT_EVERY):
            batch = file_tasks[batch_start:batch_start + CHECKPOINT_EVERY]

            # 1) Извлечение текста (локально, быстро)
            batch_chunks = []  # (file_name, chunk)
            for file_path, file_name in batch:
                ext = os.path.splitext(file_name)[1].lower()
                try:
                    chunks = extract_text_chunks(file_path, ext)
                except Exception as e:
                    print(f"  ✗ Ошибка чтения {file_name}: {e}")
                    continue

                print(f"  📄 {file_name}: чанков — {len(chunks)}")
                store.add_document_metadata(file_name, ext.replace(".", "").upper())
                all_processed_documents.append(file_name)
                batch_chunks.extend((file_name, chunk) for chunk in chunks)

            # 2) Параллельный LLM-анализ всех чанков батча
            futures = {
                executor.submit(get_ai_analysis, chunk): file_name
                for file_name, chunk in batch_chunks
            }
            done_chunks = 0
            for future in as_completed(futures):
                file_name = futures[future]
                store.merge_analysis(future.result(), file_name)
                done_chunks += 1
                print(f"  📝 проанализировано чанков: {done_chunks}/{len(futures)}", end="\r")
            if futures:
                print()

            processed = min(batch_start + CHECKPOINT_EVERY, total)

            # 3) Checkpoint
            store.save()
            elapsed = time.time() - start_time
            rate = processed / elapsed * 60
            eta_min = (total - processed) / rate if rate > 0 else 0
            print(f"💾 {processed}/{total} файлов | чанков в батче: {done_chunks} | "
                  f"{elapsed/60:.1f} мин | ~{rate:.1f} файлов/мин | осталось ~{eta_min:.0f} мин")

    # ===== КРОСС-ДОКУМЕНТНЫЕ СВЯЗИ =====
    print(f"\n{'='*50}")
    print(f"🔗 Построение кросс-документных связей")
    print(f"{'='*50}")
    cross_links = build_cross_document_links(store.nodes)

    # ===== ФИНАЛЬНОЕ СОХРАНЕНИЕ =====
    edges_df = deduplicate_edges(pd.DataFrame(store.edges, columns=EDGES_COLS))
    store.save()
    edges_df.to_csv(BASE_DIR / "edges.csv", index=False, encoding="utf-8-sig")

    elapsed = time.time() - start_time
    print(f"\n{'='*50}")
    print(f"✅ ЗАВЕРШЕНО за {elapsed/60:.1f} минут!")
    print(f"{'='*50}")
    print(f"📊 ИТОГОВАЯ СТАТИСТИКА:")
    print(f"  📌 Узлов (сущностей): {len(store.nodes)}")
    print(f"  🔗 Связей (после дедупликации): {len(edges_df)}")
    print(f"  ↔️  Двусторонних связей: {int(edges_df['bidirectional'].sum()) if 'bidirectional' in edges_df.columns else 0}")
    print(f"  📊 Параметров: {len(store.parameters)}")
    print(f"  👤 Экспертов: {len(store.experts)}")
    print(f"  📄 Документов в базе: {len(store.document_metadata)}")
    print(f"  ⚠️  Противоречий найдено: {len(store.contradictions)}")
    print(f"  🌐 Кросс-документных связей: {len(cross_links)}")
    print(f"  ✓ Документов обработано: {len(all_processed_documents)}")
    print(f"\n💾 Таблицы сохранены в: {BASE_DIR}")

if __name__ == "__main__":
    main()
